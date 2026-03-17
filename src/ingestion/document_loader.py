"""
Carga de documentos desde distintas fuentes y formatos.
Soporta: PDF, PPTX, DOCX, XLSX, TXT, Markdown.
"""
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import List, Optional

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation

from src.utils.logger import setup_logger

logger = setup_logger(__name__)


class DocumentType(str, Enum):
    """Tipos de documento soportados por el sistema."""

    PDF = "pdf"
    PPTX = "pptx"
    DOCX = "docx"
    XLSX = "xlsx"
    TXT = "txt"
    MARKDOWN = "markdown"


class SourceCategory(str, Enum):
    """Categoría de la fuente del documento."""

    INFORME_GESTION = "informe_gestion"
    ESTADO_FINANCIERO = "estado_financiero"
    NORMATIVA = "normativa"
    OTRO = "otro"


@dataclass
class LoadedDocument:
    """Representa un documento cargado listo para parseo y chunking."""

    doc_id: str
    file_name: str
    file_path: str
    doc_type: DocumentType
    source_category: SourceCategory
    country: Optional[str]          # honduras, guatemala, nicaragua
    pages_or_slides: int
    raw_text: str
    tables: List[dict] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)


class DocumentLoader:
    """
    Carga documentos de distintos formatos y retorna estructura normalizada.
    Soporta PDF, PPTX, DOCX, TXT y Markdown.
    """

    SUPPORTED_EXTENSIONS = {
        ".pdf": DocumentType.PDF,
        ".pptx": DocumentType.PPTX,
        ".docx": DocumentType.DOCX,
        ".txt": DocumentType.TXT,
        ".md": DocumentType.MARKDOWN,
    }

    def load_file(
        self,
        file_path: Path,
        source_category: SourceCategory = SourceCategory.OTRO,
        country: Optional[str] = None,
    ) -> Optional[LoadedDocument]:
        """
        Carga un archivo y retorna un LoadedDocument normalizado.

        Args:
            file_path: Ruta absoluta al archivo.
            source_category: Categoría semántica del documento.
            country: País de origen (para informes de gestión).

        Returns:
            LoadedDocument con texto extraído, o None si el formato no es soportado.
        """
        suffix = file_path.suffix.lower()
        doc_type = self.SUPPORTED_EXTENSIONS.get(suffix)

        if doc_type is None:
            logger.warning("Formato no soportado: %s", suffix)
            return None

        logger.info("Cargando documento: %s (%s)", file_path.name, doc_type.value)

        loaders = {
            DocumentType.PDF: self._load_pdf,
            DocumentType.PPTX: self._load_pptx,
            DocumentType.DOCX: self._load_docx,
            DocumentType.TXT: self._load_text,
            DocumentType.MARKDOWN: self._load_text,
        }

        raw_text, tables, pages = loaders[doc_type](file_path)

        return LoadedDocument(
            doc_id=self._generate_doc_id(file_path),
            file_name=file_path.name,
            file_path=str(file_path),
            doc_type=doc_type,
            source_category=source_category,
            country=country,
            pages_or_slides=pages,
            raw_text=raw_text,
            tables=tables,
            metadata={
                "source_category": source_category.value,
                "country": country,
                "doc_type": doc_type.value,
                "file_name": file_path.name,
            },
        )

    def load_directory(
        self,
        directory: Path,
        source_category: SourceCategory = SourceCategory.OTRO,
        country: Optional[str] = None,
    ) -> List[LoadedDocument]:
        """
        Carga todos los documentos soportados de un directorio.

        Args:
            directory: Ruta al directorio.
            source_category: Categoría semántica para todos los docs del directorio.
            country: País de origen para todos los docs del directorio.

        Returns:
            Lista de LoadedDocument cargados exitosamente.
        """
        documents = []
        for file_path in sorted(directory.iterdir()):
            if file_path.is_file() and file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS:
                doc = self.load_file(file_path, source_category, country)
                if doc is not None:
                    documents.append(doc)
        logger.info(
            "Directorio %s: %d documentos cargados", directory.name, len(documents)
        )
        return documents

    def _load_pdf(self, file_path: Path) -> tuple[str, List[dict], int]:
        """Extrae texto y tablas de un PDF usando pdfplumber."""
        all_text = []
        all_tables = []

        with pdfplumber.open(file_path) as pdf:
            pages = len(pdf.pages)
            for page_num, page in enumerate(pdf.pages, start=1):
                # Extrae texto de la página
                text = page.extract_text() or ""
                if text.strip():
                    all_text.append(f"[Página {page_num}]\n{text}")

                # Extrae tablas si las hay
                for table in page.extract_tables():
                    if table:
                        all_tables.append({
                            "page": page_num,
                            "data": table,
                        })

        return "\n\n".join(all_text), all_tables, pages

    def _load_pptx(self, file_path: Path) -> tuple[str, List[dict], int]:
        """Extrae texto y tablas de una presentación PPTX."""
        prs = Presentation(str(file_path))
        all_text = []
        all_tables = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []

            for shape in slide.shapes:
                # Extrae texto de shapes con texto
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)

                # Extrae tablas
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        table_data.append([cell.text.strip() for cell in row.cells])
                    all_tables.append({"slide": slide_num, "data": table_data})

            if slide_texts:
                all_text.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

        return "\n\n".join(all_text), all_tables, len(prs.slides)

    def _load_docx(self, file_path: Path) -> tuple[str, List[dict], int]:
        """Extrae texto y tablas de un documento DOCX."""
        doc = DocxDocument(str(file_path))
        all_text = []
        all_tables = []

        for para in doc.paragraphs:
            if para.text.strip():
                all_text.append(para.text.strip())

        for table_num, table in enumerate(doc.tables, start=1):
            table_data = []
            for row in table.rows:
                table_data.append([cell.text.strip() for cell in row.cells])
            all_tables.append({"table_num": table_num, "data": table_data})

        return "\n\n".join(all_text), all_tables, len(doc.paragraphs)

    def _load_text(self, file_path: Path) -> tuple[str, List[dict], int]:
        """Carga texto plano o Markdown."""
        text = file_path.read_text(encoding="utf-8")
        lines = len(text.splitlines())
        return text, [], lines

    def _generate_doc_id(self, file_path: Path) -> str:
        """Genera un ID único basado en el nombre del archivo."""
        import hashlib
        return hashlib.sha1(str(file_path).encode()).hexdigest()[:12]
